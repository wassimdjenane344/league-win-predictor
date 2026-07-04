"""Generate the project presentation (English) as a .pptx file.

Run:  python presentation/build_deck.py
Output: presentation/League_Win_Predictor.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- palette (matches the LoL-themed app) ---
NAVY = RGBColor(0x0A, 0x14, 0x28)
PANEL = RGBColor(0x0F, 0x21, 0x40)
PANEL2 = RGBColor(0x14, 0x2A, 0x50)
GOLD = RGBColor(0xC8, 0xAA, 0x6E)
TEXT = RGBColor(0xF0, 0xE6, 0xD2)
MUTED = RGBColor(0x9A, 0xB0, 0xD0)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
LINE = RGBColor(0x24, 0x49, 0x6F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

FONT = "Segoe UI"


def slide():
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    return s


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, bullet=False, space=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    return p


def title(s, text, kicker=None):
    if kicker:
        tf = textbox(s, 0.6, 0.35, 12.1, 0.5)
        para(tf, kicker.upper(), size=13, color=GOLD, bold=True, first=True)
    tf = textbox(s, 0.6, 0.75, 12.1, 1.0)
    para(tf, text, size=30, color=TEXT, bold=True, first=True)
    # gold underline
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.55), Inches(2.2), Pt(3))
    ln.fill.solid()
    ln.fill.fore_color.rgb = GOLD
    ln.line.fill.background()


def box(s, x, y, w, h, lines, fill=PANEL, border=GOLD, title_color=GOLD, body_color=TEXT, title_size=15, body_size=11):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, (txt, is_title) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.bold = is_title
        r.font.size = Pt(title_size if is_title else body_size)
        r.font.color.rgb = title_color if is_title else body_color
    return shp


def arrow(s, x, y, w, h=0.3, color=GOLD, shape=MSO_SHAPE.RIGHT_ARROW):
    a = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def connector(s, x1, y1, x2, y2, color=GOLD, width=1.5, dash=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def panel(s, x, y, w, h, fill=PANEL, border=LINE):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def bullets(s, items, x=0.9, y=1.9, w=11.5, size=18, gap=10, color=TEXT):
    tf = textbox(s, x, y, w, 5.0)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, col, bold = it
        else:
            txt, col, bold = it, color, False
        para(tf, txt, size=size, color=col, bold=bold, bullet=True, space=gap, first=(i == 0))


# =========================================================================
# SLIDE 1 — Title
# =========================================================================
s = slide()
# subtle top band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
band.fill.solid()
band.fill.fore_color.rgb = GOLD
band.line.fill.background()
tf = textbox(s, 0.8, 2.2, 11.7, 1.2)
para(tf, "League Win Predictor", size=48, color=GOLD, bold=True, first=True)
tf = textbox(s, 0.8, 3.35, 11.7, 1.0)
para(tf, "An End-to-End MLOps Project", size=26, color=TEXT, first=True)
tf = textbox(s, 0.8, 4.3, 11.7, 1.2)
para(tf, "Predicting a League of Legends win from the first 10 minutes —", size=17, color=MUTED, first=True)
para(tf, "and, more importantly, running that AI model in production the professional way.", size=17, color=MUTED)
tf = textbox(s, 0.8, 6.4, 11.7, 0.8)
para(tf, "DevOps & MLOps — Final Project   •   Team of 3   •   Live: lol-frontend-scaw.onrender.com",
     size=13, color=GOLD, first=True)

# =========================================================================
# SLIDE 2 — The goal
# =========================================================================
s = slide()
title(s, "What is this project really about?", kicker="The goal")
bullets(s, [
    ("Building an AI model that works on your laptop is the easy part.", TEXT, True),
    ("The real challenge — and what this project is graded on — is running that model", MUTED, False),
    ("in production: reliably, automatically, traceably, and monitored.", MUTED, False),
], y=1.9, gap=6)
panel(s, 0.9, 3.4, 11.5, 3.2)
tf = textbox(s, 1.3, 3.7, 10.7, 2.7, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "The full MLOps lifecycle we had to build:", size=17, color=GOLD, bold=True, first=True, space=10)
for t in [
    "Data versioning  —  know exactly which data trained which model",
    "Model training, tracking & a model registry",
    "CI/CD with quality gates  —  nothing ships unless it passes",
    "Reproducible deployments to the cloud",
    "Live monitoring of the production service",
]:
    para(tf, t, size=15, color=TEXT, bullet=True, space=6)
tf = textbox(s, 0.9, 6.75, 11.5, 0.6)
para(tf, "The game (LoL) is just the pretext. The engineering around it is the point.",
     size=14, color=MUTED, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
# SLIDE 3 — What the app does
# =========================================================================
s = slide()
title(s, "The app in one look", kicker="The product")
bullets(s, [
    ("A web app that predicts the blue team's win probability at minute 10", TEXT, True),
    ("of a ranked League of Legends game, from stats like gold lead, kills, dragons…", MUTED, False),
], y=1.9, gap=6)
# mock result bar
panel(s, 0.9, 3.2, 11.5, 2.9)
tf = textbox(s, 1.3, 3.45, 10.7, 0.6)
para(tf, "You enter the match stats  →  you get a prediction:", size=16, color=TEXT, first=True)
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(4.2), Inches(5.8), Inches(0.6))
b.fill.solid(); b.fill.fore_color.rgb = BLUE; b.line.fill.background()
tfb = b.text_frame; tfb.word_wrap = True
pb = tfb.paragraphs[0]; pb.alignment = PP_ALIGN.CENTER
rb = pb.add_run(); rb.text = "Blue 55%"; rb.font.bold = True; rb.font.size = Pt(14); rb.font.color.rgb = RGBColor(255,255,255); rb.font.name = FONT
r2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(4.2), Inches(4.7), Inches(0.6))
r2.fill.solid(); r2.fill.fore_color.rgb = RED; r2.line.fill.background()
tfr = r2.text_frame
pr = tfr.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
rr = pr.add_run(); rr.text = "Red 45%"; rr.font.bold = True; rr.font.size = Pt(14); rr.font.color.rgb = RGBColor(255,255,255); rr.font.name = FONT
tf = textbox(s, 1.3, 5.05, 10.7, 0.9)
para(tf, "Every prediction is traceable: it shows the model version, the git commit and the",
     size=12, color=MUTED, first=True)
para(tf, "data version that produced it. The model runs on a real server, served from the registry.", size=12, color=MUTED)

# =========================================================================
# SLIDE 4 — Restaurant analogy
# =========================================================================
s = slide()
title(s, "The mental model: a restaurant", kicker="How to think about it")
tf = textbox(s, 0.9, 1.75, 11.5, 0.6)
para(tf, "Cooking a great dish at home = your model working on your laptop.  Opening a "
         "restaurant that serves it to hundreds of customers every day = this project.",
     size=15, color=MUTED, first=True)
pairs = [
    ("The model", "The recipe / the dish"),
    ("Data + DVC", "The ingredients — and knowing exactly which batch was used"),
    ("MLflow", "The chef's logbook + the shelf holding the official recipe"),
    ("Tests", "Quality control — taste before serving"),
    ("Quality gate", "The chef tastes: good → served; bad → stays in the kitchen"),
    ("CI/CD", "The automatic conveyor: kitchen → QC → dining room"),
    ("Git branches", "Draft → test kitchen → dress-rehearsal → dining room"),
    ("Docker", "The standard container — same everywhere"),
    ("Cloud (Render)", "The restaurant, open to the public, with an address"),
    ("Prometheus/Grafana", "The cameras watching the dining room, live"),
]
y = 2.5
for i, (left, right) in enumerate(pairs):
    yy = 2.45 + (i % 5) * 0.92
    xx = 0.9 if i < 5 else 6.95
    b = panel(s, xx, yy, 5.6, 0.82, fill=PANEL if i % 2 == 0 else PANEL2)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = left + "  —  "; r.font.bold = True; r.font.size = Pt(12.5); r.font.color.rgb = GOLD; r.font.name = FONT
    r2 = p.add_run(); r2.text = right; r2.font.size = Pt(11.5); r2.font.color.rgb = TEXT; r2.font.name = FONT

# =========================================================================
# SLIDE 5 — Architecture (THE BIG DIAGRAM)
# =========================================================================
s = slide()
title(s, "Architecture — the whole process", kicker="The big picture")

row_y = 3.05
bh = 1.2
box(s, 0.30, row_y, 2.15, bh, [("Developer", True), ("your laptop", False), ("write code + tests", False)], fill=PANEL2)
arrow(s, 2.52, row_y + 0.45, 0.33)
box(s, 2.95, row_y, 2.15, bh, [("GitHub", True), ("code + 4 branches", False), ("feature→dev→staging→main", False)], fill=PANEL2)
arrow(s, 5.17, row_y + 0.45, 0.33)
box(s, 5.60, row_y, 2.35, bh, [("GitHub Actions", True), ("3 CI/CD pipelines", False), ("test • build • gate", False)], fill=PANEL, border=BLUE, title_color=BLUE)
arrow(s, 8.00, row_y + 0.45, 0.33)
box(s, 8.45, row_y, 2.35, bh, [("Render (cloud)", True), ("Backend (Flask+model)", False), ("Frontend (Next.js)", False)], fill=PANEL, border=GREEN, title_color=GREEN)
arrow(s, 10.85, row_y + 0.45, 0.30)
box(s, 11.20, row_y, 1.90, bh, [("Users", True), ("public URL", False)], fill=PANEL2)

# DagsHub above
box(s, 6.55, 1.15, 3.05, 1.15, [("DagsHub", True), ("MLflow registry + DVC data", False), ("experiments • models • dataset", False)], fill=PANEL, border=GOLD)
connector(s, 6.9, row_y, 7.6, 2.30)          # CI/CD -> DagsHub
connector(s, 9.0, 2.30, 9.6, row_y)          # DagsHub -> Render

# Monitoring below
box(s, 8.45, 4.95, 2.35, 1.1, [("Prometheus + Grafana", True), ("live dashboards", False)], fill=PANEL, border=RED, title_color=RED)
connector(s, 9.6, 4.95, 9.6, row_y + bh)     # monitoring -> Render backend

# small labels
def lbl(x, y, text, color=MUTED, size=9.5, w=2.6):
    tf = textbox(s, x, y, w, 0.35)
    para(tf, text, size=size, color=color, align=PP_ALIGN.CENTER, first=True)

lbl(6.2, 2.35, "train • register • promote", size=9)
lbl(8.7, 2.35, "serve Production model", size=9)
lbl(9.7, 4.55, "scrape /metrics", size=9, w=2.0)

# caption
panel(s, 0.6, 6.35, 12.1, 0.95, fill=PANEL)
tf = textbox(s, 0.9, 6.45, 11.6, 0.8, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "You push code → CI/CD tests it → it trains & promotes the model on DagsHub only if the "
         "quality gate passes → it deploys to Render → users get predictions.", size=12.5, color=TEXT,
     align=PP_ALIGN.CENTER, first=True, space=2)
para(tf, "Prometheus & Grafana monitor the live production backend.", size=12.5, color=MUTED,
     align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 6 — DevOps vs MLOps
# =========================================================================
s = slide()
title(s, "DevOps and MLOps — how they split", kicker="The two halves")
# two columns
box_dev = panel(s, 0.9, 1.9, 5.6, 4.6, fill=PANEL)
tf = box_dev.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
para(tf, "DevOps", size=22, color=BLUE, bold=True, first=True, space=4)
para(tf, "automating build, test & delivery", size=12, color=MUTED, space=12)
for t in ["Tests (unit + integration + e2e)", "CI/CD (the 3 GitHub pipelines)",
          "Git hooks (checks before push)", "Docker (packaging)",
          "Auto-deploy to Render", "Monitoring (Prometheus / Grafana)"]:
    para(tf, t, size=14, color=TEXT, bullet=True, space=8)

box_ml = panel(s, 6.85, 1.9, 5.6, 4.6, fill=PANEL)
tf = box_ml.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.2)
para(tf, "MLOps", size=22, color=GOLD, bold=True, first=True, space=4)
para(tf, "the same discipline, applied to data & models", size=12, color=MUTED, space=12)
for t in ["DVC (versioning the data)", "MLflow (tracking + model registry)",
          "The quality gate + model promotion", "Serving the model from the registry",
          "Traceability: data version + git commit", "= DevOps + the AI-specific pieces"]:
    para(tf, t, size=14, color=TEXT, bullet=True, space=8)
tf = textbox(s, 0.9, 6.7, 11.5, 0.6)
para(tf, "MLOps = DevOps extended to the model and the data, not just the code.",
     size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
# SLIDE 7 — Tech stack
# =========================================================================
s = slide()
title(s, "The technology stack", kicker="What we used")
stack = [
    ("Frontend", "Next.js (React)"),
    ("Backend / API", "Python + Flask"),
    ("ML", "scikit-learn"),
    ("Data versioning", "DVC"),
    ("Experiments & registry", "MLflow + DagsHub"),
    ("Packaging", "Docker + docker-compose"),
    ("CI/CD", "GitHub Actions"),
    ("Monitoring", "Prometheus + Grafana"),
    ("Cloud hosting", "Render"),
]
for i, (k, v) in enumerate(stack):
    xx = 0.9 + (i % 3) * 3.95
    yy = 2.2 + (i // 3) * 1.35
    b = panel(s, xx, yy, 3.7, 1.15, fill=PANEL if (i // 3) % 2 == 0 else PANEL2)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    para(tf, k, size=12, color=MUTED, first=True, space=2)
    para(tf, v, size=16, color=GOLD, bold=True)

# =========================================================================
# SLIDE 8 — Data + DVC
# =========================================================================
s = slide()
title(s, "Data versioning with DVC", kicker="MLOps — data")
bullets(s, [
    ("The real Kaggle dataset: 9,879 Diamond-ranked LoL games (first 10 min).", TEXT, True),
    ("Git is great for code, but not for large data files. DVC versions the data:", MUTED, False),
    ("Git stores a tiny pointer; the actual dataset lives in remote storage (DagsHub).", MUTED, False),
    ("Every training run records which data version + which git commit produced it.", MUTED, False),
], y=1.9, gap=10)
panel(s, 0.9, 4.6, 11.5, 1.9, fill=PANEL)
tf = textbox(s, 1.2, 4.8, 11.0, 1.6, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Why it matters:", size=16, color=GOLD, bold=True, first=True, space=6)
para(tf, "Reproducibility. Six months later, we can rebuild the exact model from the exact "
         "data — and prove which dataset any deployed model was trained on.", size=14, color=TEXT)

# =========================================================================
# SLIDE 9 — Model + MLflow
# =========================================================================
s = slide()
title(s, "Model tracking & registry with MLflow", kicker="MLOps — models")
bullets(s, [
    ("Every training run logs its parameters, metrics, data version and git commit.", TEXT, True),
    ("The Model Registry is the single source of truth for what is deployed.", MUTED, False),
    ("Models move through stages:", MUTED, False),
], y=1.9, gap=10)
# stage boxes
box(s, 1.1, 3.9, 3.2, 1.1, [("Staging", True), ("candidate model", False), ("under validation", False)], fill=PANEL, border=GOLD, title_color=GOLD, title_size=17)
arrow(s, 4.45, 4.3, 0.5)
box(s, 5.15, 3.9, 3.2, 1.1, [("Production", True), ("the model users get", False), ("served live", False)], fill=PANEL, border=GREEN, title_color=GREEN, title_size=17)
box(s, 9.2, 3.9, 3.1, 1.1, [("Archived", True), ("previous versions", False)], fill=PANEL2, border=LINE, title_color=MUTED, title_size=17)
tf = textbox(s, 0.9, 5.5, 11.5, 1.2)
para(tf, "Production serves the \"Production\" model ONLY. A new model becomes Production "
         "only after passing the quality gate (next slide).", size=15, color=TEXT, first=True)

# =========================================================================
# SLIDE 10 — Git branching
# =========================================================================
s = slide()
title(s, "Strict Git branching model", kicker="DevOps — workflow")
stages = [
    ("feature/*", "where you develop", MUTED),
    ("dev", "integration", BLUE),
    ("staging", "pre-production validation", GOLD),
    ("main", "production (live)", GREEN),
]
for i, (name, desc, col) in enumerate(stages):
    xx = 0.6 + i * 3.15
    box(s, xx, 2.7, 2.7, 1.3, [(name, True), (desc, False)], fill=PANEL, border=col, title_color=col, title_size=18, body_size=12)
    if i < 3:
        arrow(s, xx + 2.72, 3.25, 0.38)
tf = textbox(s, 0.9, 4.6, 11.5, 1.8)
para(tf, "Each arrow triggers one of our CI/CD pipelines automatically:", size=16, color=GOLD, bold=True, first=True, space=8)
para(tf, "Code only reaches production by travelling feature → dev → staging → main,", size=15, color=TEXT)
para(tf, "passing automated checks at every step. Nothing is deployed by hand.", size=15, color=TEXT)

# =========================================================================
# SLIDE 11 — The 3 pipelines
# =========================================================================
s = slide()
title(s, "The 3 CI/CD pipelines (we wrote them)", kicker="DevOps — automation")
tf = textbox(s, 0.9, 1.7, 11.5, 0.5)
para(tf, "GitHub provides the robot (GitHub Actions); we wrote its instructions in .github/workflows/.",
     size=13, color=MUTED, first=True)
data = [
    ("PR → dev", BLUE, ["Lint the code", "Run unit tests", "Run integration tests", "Build Docker images"]),
    ("dev → staging", GOLD, ["Full test suite + e2e", "Train a candidate model", "Register it as 'Staging'"]),
    ("staging → main", GREEN, ["Quality gate check", "Promote model to Production", "Auto-deploy to Render"]),
]
for i, (name, col, steps) in enumerate(data):
    xx = 0.9 + i * 3.95
    b = panel(s, xx, 2.35, 3.7, 4.0, fill=PANEL)
    tf = b.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.18)
    para(tf, name, size=17, color=col, bold=True, first=True, space=10)
    for st in steps:
        para(tf, st, size=13.5, color=TEXT, bullet=True, space=8)
tf = textbox(s, 0.9, 6.55, 11.5, 0.6)
para(tf, "All three are currently green in CI.", size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
# SLIDE 12 — Quality gate (core)
# =========================================================================
s = slide()
title(s, "The quality gate & model promotion", kicker="The core requirement")
tf = textbox(s, 0.9, 1.75, 11.5, 0.6)
para(tf, "When we push to main, the pipeline decides whether the new model is allowed into production:",
     size=15, color=MUTED, first=True)
box(s, 4.9, 2.5, 3.5, 0.95, [("Train candidate", True), ("evaluate accuracy", False)], fill=PANEL, title_size=15)
arrow(s, 6.4, 3.5, 0.35, 0.35, shape=MSO_SHAPE.DOWN_ARROW)
box(s, 4.4, 3.95, 4.5, 0.95, [("Quality gate: accuracy ≥ 70% ?", True)], fill=PANEL, border=GOLD, title_size=15)
# PASS
arrow(s, 3.7, 4.45, 0.6, 0.32, color=GREEN, shape=MSO_SHAPE.LEFT_ARROW)
box(s, 0.7, 5.1, 3.4, 1.1, [("PASS", True), ("promote to Production", False), ("+ deploy to Render", False)], fill=PANEL, border=GREEN, title_color=GREEN, title_size=16)
connector(s, 5.0, 4.9, 2.4, 5.1, color=GREEN, width=1.5)
# FAIL
arrow(s, 9.0, 4.45, 0.6, 0.32, color=RED)
box(s, 9.5, 5.1, 3.3, 1.1, [("FAIL", True), ("stays in Staging", False), ("production unchanged", False)], fill=PANEL, border=RED, title_color=RED, title_size=16)
connector(s, 8.3, 4.9, 11.1, 5.1, color=RED, width=1.5)
tf = textbox(s, 0.9, 6.5, 11.5, 0.6)
para(tf, "A bad model can never reach users. This is what makes updating the AI safe.",
     size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
# SLIDE 13 — Testing
# =========================================================================
s = slide()
title(s, "Testing — the pyramid", kicker="DevOps — quality")
rows = [
    ("End-to-end (1)", "A real browser fills the form and gets a real prediction", RED, 4.0),
    ("Integration (2)", "The Flask API + the loaded model work together (/health, /predict)", GOLD, 7.0),
    ("Unit (3+)", "Small pure functions in isolation (feature builder, quality gate…)", BLUE, 10.0),
]
y = 2.4
for name, desc, col, w in rows:
    xx = (13.333 - w) / 2
    b = box(s, xx, y, w, 1.15, [(name, True)], fill=PANEL, border=col, title_color=col, title_size=17)
    y += 1.35
    tf = textbox(s, 1.0, y - 0.15, 11.3, 0.5)
    para(tf, desc, size=12.5, color=MUTED, align=PP_ALIGN.CENTER, first=True)
    y += 0.25
tf = textbox(s, 0.9, 6.55, 11.5, 0.6)
para(tf, "All of these run automatically in CI on every pull request.",
     size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
# SLIDE 14 — Docker
# =========================================================================
s = slide()
title(s, "Docker — the same everywhere", kicker="DevOps — packaging")
bullets(s, [
    ("Each service (backend, frontend, trainer) has a Dockerfile: a recipe that", TEXT, True),
    ("packages the app with everything it needs — no more \"works on my machine\".", MUTED, False),
    ("Render builds and runs these exact images — your live app runs inside Docker right now.", MUTED, False),
    ("docker-compose runs the whole stack locally with one command.", MUTED, False),
], y=1.9, gap=10)
panel(s, 0.9, 4.5, 11.5, 2.0, fill=PANEL)
tf = textbox(s, 1.2, 4.7, 11.0, 1.7, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "backend/Dockerfile, in plain English:", size=14, color=GOLD, bold=True, first=True, space=8)
para(tf, "\"Start from Python → install my libraries → copy my code → launch the web server.\"",
     size=15, color=TEXT, space=4)
para(tf, "The result is a self-contained box that runs identically on any machine.", size=13, color=MUTED)

# =========================================================================
# SLIDE 15 — Monitoring
# =========================================================================
s = slide()
title(s, "Monitoring the live production", kicker="DevOps — observability")
bullets(s, [
    ("The backend exposes a /metrics endpoint.", TEXT, True),
    ("Prometheus scrapes it; Grafana turns it into a live dashboard.", MUTED, False),
    ("It runs against the LIVE production backend, not just locally.", MUTED, False),
], y=1.9, gap=8)
metrics = ["Number of predictions", "Prediction latency", "Failed requests", "Backend health / uptime"]
for i, m in enumerate(metrics):
    xx = 0.9 + (i % 2) * 5.95
    yy = 3.7 + (i // 2) * 1.25
    b = panel(s, xx, yy, 5.6, 1.05, fill=PANEL if i % 2 == 0 else PANEL2)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.25)
    para(tf, m, size=16, color=GOLD, bold=True, first=True)
tf = textbox(s, 0.9, 6.5, 11.5, 0.6)
para(tf, "docker compose -f docker-compose.monitoring.yml up   →   Grafana at localhost:3001",
     size=13, color=MUTED, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
# SLIDE 16 — Live demo
# =========================================================================
s = slide()
title(s, "Live demo", kicker="Let's show it")
steps = [
    "1.  Open the app and make a prediction (lol-frontend-scaw.onrender.com)",
    "2.  Show the code + README on GitHub",
    "3.  Show the 3 green pipelines in the Actions tab",
    "4.  Show the model registry on DagsHub (a model in Production)",
    "5.  Show the Grafana monitoring dashboard",
]
bullets(s, [(st, TEXT, False) for st in steps], y=2.2, gap=16, size=18)
panel(s, 0.9, 5.7, 11.5, 1.0, fill=PANEL)
tf = textbox(s, 1.2, 5.85, 11.0, 0.8, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Tip: the free server may take ~30s to wake up on the first request — that's normal.",
     size=13, color=MUTED, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
# SLIDE 17 — Conclusion
# =========================================================================
s = slide()
title(s, "What we proved", kicker="Conclusion")
bullets(s, [
    ("A real AI model, live in production, served from a model registry.", TEXT, True),
    ("A fully automated pipeline: test → train → quality gate → deploy.", TEXT, True),
    ("Every model traceable to its data version and git commit.", TEXT, True),
    ("Monitored in real time.", TEXT, True),
], y=1.9, gap=12, size=17)
panel(s, 0.9, 4.6, 11.5, 1.9, fill=PANEL)
tf = textbox(s, 1.2, 4.8, 11.0, 1.6, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "The value:", size=15, color=GOLD, bold=True, first=True, space=6)
para(tf, "In real companies, the hard part of AI isn't building a model — it's operating it in "
         "production safely and repeatably. That is the MLOps engineer's job, and it's exactly "
         "what this project demonstrates end to end.", size=14, color=TEXT)
tf = textbox(s, 0.9, 6.7, 11.5, 0.6)
para(tf, "Thank you  —  Questions?", size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER, first=True)

# =========================================================================
out = Path(__file__).with_name("League_Win_Predictor.pptx")
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides")
